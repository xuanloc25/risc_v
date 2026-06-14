# -*- coding: utf-8 -*-
"""Dựng deck phản biện KLTN trên theme mẫu Khoa KTMT-UIT.
Kế thừa master/theme/layout từ template, thay 15 slide mẫu bằng 22 slide nội dung,
nhúng ghi chú (speaker notes) và dashboard kiểm chứng.
"""
import copy
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

BASE = r"D:\STUDY\UIT\DoAn1\web\assembler\docs"
TEMPLATE = BASE + r"\ref\KTMT_KLTN_Phu luc 6_Mau bao cao bao ve.pptx"
OUT = BASE + r"\kltn\slides\KLTN_phan_bien.pptx"
DASH = BASE + r"\kltn\figures\fig_5_1_verification_dashboard.png"

prs = Presentation(TEMPLATE)
SW, SH = prs.slide_width, prs.slide_height

# Layout indices (từ kiểm tra template)
L_TITLE = prs.slide_layouts[0]       # Title Slide
L_SECTION = prs.slide_layouts[1]     # Section Header
L_CONTENT = {1: prs.slide_layouts[2], 2: prs.slide_layouts[6], 3: prs.slide_layouts[5],
             4: prs.slide_layouts[4], 5: prs.slide_layouts[3]}  # 1_..5_ Title and Content
L_TITLE_ONLY = prs.slide_layouts[9]

# --- Xóa toàn bộ slide mẫu hiện có (drop cả relationship để không trùng part) ---
sldIdLst = prs.slides._sldIdLst
for sid in list(sldIdLst):
    rId = sid.get(qn('r:id'))
    if rId:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
    sldIdLst.remove(sid)

def get_ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None

def set_title(slide, text):
    if slide.shapes.title is not None:
        slide.shapes.title.text = text

def fill_bullets(tf, bullets):
    tf.word_wrap = True
    first = True
    for b in bullets:
        lvl = 1 if b.startswith("  ") else 0
        b = b.strip()
        if first:
            tf.text = b
            tf.paragraphs[0].level = lvl
            first = False
        else:
            p = tf.add_paragraph(); p.text = b; p.level = lvl

def set_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes

def add_content(section, title, bullets, notes):
    slide = prs.slides.add_slide(L_CONTENT[section])
    set_title(slide, title)
    ph = get_ph(slide, 1)
    if ph is not None:
        fill_bullets(ph.text_frame, bullets)
    set_notes(slide, notes)
    return slide

def add_divider(title, number, notes=""):
    slide = prs.slides.add_slide(L_SECTION)
    set_title(slide, title)
    body = get_ph(slide, 1)
    if body is not None:
        body.text_frame.text = number
    if notes:
        set_notes(slide, notes)
    return slide

def _title_bottom(layout):
    """Đáy của title placeholder trong layout (EMU); fallback nếu thiếu."""
    for ph in layout.placeholders:
        if ph.placeholder_format.idx == 0 and ph.top is not None and ph.height is not None:
            return ph.top + ph.height
    return Emu(1500000)

def add_image_slide(section, title, bullets, notes, image):
    """Title-only + textbox bullets (ngay dưới tiêu đề) + ảnh (phần dưới)."""
    slide = prs.slides.add_slide(L_TITLE_ONLY)
    set_title(slide, title)
    tb_top = _title_bottom(L_TITLE_ONLY) + Emu(140000)
    tb_h = Emu(1350000)
    tb = slide.shapes.add_textbox(Emu(457200), tb_top, SW - Emu(914400), tb_h)
    fill_bullets(tb.text_frame, bullets)
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(13)
    # ảnh: căn giữa, nằm dưới textbox, chừa lề đáy
    from PIL import Image
    try:
        iw, ih = Image.open(image).size
    except Exception:
        iw, ih = 1960, 1184
    img_top = tb_top + tb_h + Emu(120000)
    avail_w = SW - Emu(1828800)            # lề ~1 inch mỗi bên
    avail_h = SH - img_top - Emu(450000)   # chừa lề đáy
    ratio = min(avail_w / iw, avail_h / ih)
    w = int(iw * ratio); h = int(ih * ratio)
    left = int((SW - w) / 2)
    slide.shapes.add_picture(image, left, int(img_top), width=w, height=h)
    set_notes(slide, notes)
    return slide

def add_title_slide(title, lines, notes):
    slide = prs.slides.add_slide(L_TITLE)
    set_title(slide, title)
    body = get_ph(slide, 14) or get_ph(slide, 15) or get_ph(slide, 1)
    if body is not None:
        fill_bullets(body.text_frame, lines)
    set_notes(slide, notes)
    return slide

# ============================ NỘI DUNG ============================

# 1. Bìa
add_title_slide(
    "[TÊN KHÓA LUẬN] — Trình mô phỏng SoC RISC-V trên nền web",
    ["Khóa luận tốt nghiệp — ngành Kỹ thuật Máy tính, Trường ĐH CNTT (UIT)",
     "GVHD: [điền tên GVHD]",
     "SV1: [Họ tên] — [MSSV]   |   SV2: [Họ tên] — [MSSV]",
     "BÁO CÁO PHẢN BIỆN — Tháng 06/2026"],
    "Em chào thầy/cô. Em là [tên] cùng [tên] thực hiện khóa luận [tên đề tài] dưới hướng dẫn của [GVHD]. "
    "Hôm nay nhóm em trình bày phần phản biện. Đề tài xây dựng trình mô phỏng SoC RISC-V chạy trực tiếp "
    "trên trình duyệt, phục vụ giảng dạy Kiến trúc máy tính. CẦN ĐIỀN: tên đề tài, GVHD, họ tên + MSSV.")

# 2. Nội dung báo cáo
add_content(1, "Nội dung báo cáo",
    ["01 — Tổng quan đề tài nghiên cứu",
     "02 — Giải pháp đề xuất của đề tài",
     "03 — Kết quả của giải pháp đề xuất",
     "04 — Kết luận và hướng phát triển",
     "Tài liệu tham khảo · 05 — Hỏi & Đáp · Phụ lục"],
    "Báo cáo gồm 4 phần theo mẫu Khoa: tổng quan; giải pháp; kết quả; kết luận–hướng phát triển. "
    "Cuối cùng là hỏi đáp và phụ lục. Nói nhanh, đây là slide bản lề.")

# ---- PHẦN 01 ----
add_divider("TỔNG QUAN ĐỀ TÀI NGHIÊN CỨU", "PHẦN 01")

add_content(1, "Bối cảnh & vấn đề cần giải quyết",
    ["Học Kiến trúc máy tính cần quan sát luồng dữ liệu xuyên suốt SoC: CPU, bộ nhớ, cache, MMU, bus, DMA, ngoại vi",
     "Công cụ phổ biến (RARS, Ripes) mạnh về CPU/pipeline nhưng ít trực quan hóa cả SoC (MMU, cache phân cấp, TileLink, DMA, ngoại vi MMIO)",
     "Mô phỏng SoC đầy đủ (Spike, gem5) cài đặt phức tạp, không chạy trực tiếp trên trình duyệt",
     "Vấn đề: thiếu công cụ chạy ngay trên web, trực quan toàn SoC, phục vụ dạy/học"],
    "Khi học KTMT, SV cần THẤY dữ liệu di chuyển trong SoC: từ CPU qua MMU, cache, lên bus, tới bộ nhớ và ngoại vi. "
    "RARS/Ripes tốt cho CPU/pipeline nhưng ít trực quan toàn SoC. Spike/gem5 mạnh nhưng cài đặt phức tạp, không chạy "
    "trên trình duyệt. Vì vậy cần một công cụ chạy ngay trên web, trực quan toàn SoC.")

add_content(1, "Khảo sát công cụ hiện có & Tính mới",
    ["So sánh RARS / Ripes / QtRvSim / Spike: phạm vi mô phỏng, mức trực quan, khả năng chạy web (bảng chi tiết ở Phụ lục)",
     "Kế thừa & khác biệt với KLTN trước (Nguyễn Gia Bảo Ngọc 2025 — Phát triển trình mô phỏng SoC)",
     "Tính mới: chạy 100% trên trình duyệt, không cài đặt",
     "Trực quan hóa datapath SỐNG: sơ đồ bus động, log đa-module, bảng TLB/cache thời gian thực",
     "Tích hợp assembler + CPU + SoC trong một công cụ thống nhất"],
    "Nhóm khảo sát RARS, Ripes, QtRvSim, Spike (bảng so sánh ở phụ lục). Đề tài kế thừa, phát triển từ KLTN 2025 của "
    "bạn Bảo Ngọc. Điểm mới: chạy 100% trên trình duyệt; trực quan datapath SỐNG (sơ đồ bus động, log đa-module, bảng "
    "TLB/cache thời gian thực); tích hợp assembler+CPU+SoC. LƯU Ý: nếu hỏi khác KLTN 2025 chỗ nào, nêu phần mình thực sự "
    "làm (TileLink UL/UH + bridge, DMA descriptor/backpressure, MMU+syscall, kiểm chứng 3 lớp). Đừng overclaim.")

add_content(1, "Mục tiêu & Phạm vi (nêu trung thực)",
    ["Mục tiêu: assembler + CPU RV32IMF + MMU + cache L1/L2 + bus TileLink UL/UH + DMA + ngoại vi, chạy web & trực quan hóa",
     "Phạm vi trung thực: RV32IMF là TẬP CON; chưa có CSR/đặc quyền/ngắt (ngoại vi polling); fence chỉ mã hóa; FPU chưa đủ FCSR/fflags/frm",
     "TileLink chỉ kênh A/D mức giao dịch (UL+UH), không B/C/E, không coherence; bus một-giao-dịch",
     "AMOADD.W là lệnh RV32A duy nhất (ngoài tên đề tài); MMU tối giản 4KB fully-associative; CAN ở mức frame (tạm gác)"],
    "ĐÂY LÀ SLIDE QUAN TRỌNG NHẤT PHẦN 1 — chủ động rào phạm vi. Nói: nhóm xin nêu rõ phạm vi trung thực ngay từ đầu. "
    "RV32IMF là tập con; chưa có CSR/đặc quyền/ngắt nên ngoại vi giao tiếp bằng polling; fence mã hóa nhưng chưa thực thi; "
    "FPU chưa đủ FCSR/fflags/frm; TileLink chỉ A/D mức giao dịch; AMOADD.W là RV32A duy nhất giữ theo thống nhất với GVHD; "
    "MMU 4KB fully-associative; CAN mức frame, tạm gác. NÊU TRƯỚC để mọi câu hỏi 'đã làm X chưa' được trả lời sẵn.")

# ---- PHẦN 02 ----
add_divider("GIẢI PHÁP ĐỀ XUẤT CỦA ĐỀ TÀI", "PHẦN 02")

add_content(2, "Kiến trúc tổng thể SoC",
    ["Chuỗi datapath: CPU → MMU → Cache L1/L2 → TileLink-UH → Main Memory; nhánh ngoại vi qua cầu nối UH→UL",
     "Memory map: RAM; LED 0xFF000000; UART 0x10000000; mouse 0xFF100000; keyboard 0xFFFF0000; DMA regs 0xFFED0000",
     "Mô hình mô phỏng theo chu kỳ (tick); đồng bộ thiết bị bằng thăm dò (polling) vì chưa có ngắt",
     "(Chèn hình sơ đồ khối SoC: figures/fig_3_2_soc_simulator_architecture.svg)"],
    "Đây là kiến trúc tổng thể. Một truy cập CPU đi qua MMU dịch địa chỉ, xuống cache L1 rồi L2, ra TileLink-UH tới bộ "
    "nhớ chính. Truy cập ngoại vi rẽ qua cầu nối UH→UL. Bản đồ bộ nhớ: LED 0xFF000000, UART 0x10000000... Mô hình chạy "
    "theo chu kỳ; chưa có ngắt nên đồng bộ bằng polling. NẾU HỎI vì sao tách UL/UH: UH cho đường bộ nhớ (burst), UL cho "
    "ngoại vi đơn giản, bridge nối hai miền. CHÈN HÌNH fig_3_2.")

add_content(2, "CPU RV32IMF & Assembler",
    ["Assembler two-pass: xử lý nhãn (tiến/lùi), pseudo-instruction (li, la, call, ret), directive dữ liệu (.word/.byte/.asciiz)",
     "Bảng opcode DÙNG CHUNG (isa.js) cho assembler + CPU + tô màu CodeMirror → nhất quán 100%",
     "CPU: chu trình fetch–decode–execute; hỗ trợ nhóm I/M/F; máy trạng thái CPU–bộ nhớ",
     "(Chèn hình: fig_4_1_two_pass_assembler.svg, fig_4_3_rv32imf_cpu_core.svg)"],
    "Assembler hai lượt: lượt 1 thu thập nhãn, lượt 2 sinh mã (xử lý được nhãn tiến). Mở rộng pseudo li/la/call/ret. "
    "Điểm tâm đắc: bảng opcode dùng chung trong isa.js cho cả assembler, CPU và tô màu cú pháp → nhất quán 100%. CPU chạy "
    "fetch–decode–execute, hỗ trợ I/M/F. NẾU HỎI làm sao chắc assembler đúng: đối chiếu GNU binutils, trình bày ở phần kết quả.")

add_content(2, "MMU & Cache phân cấp",
    ["MMU: dịch VA→PA, identity fallback khi chưa map; TLB 8 entry fully-associative LRU; kiểm tra quyền R/W/X/cacheable; map/unmap qua syscall",
     "Cache: L1I/L1D 16 set × 4 way, L2 64 set × 4 way, block 64B, write-through, LRU",
     "Vùng MMIO non-cacheable được bypass để ghi thiết bị có hiệu lực ngay",
     "(Chèn hình: fig_4_5_mmu_translation_path.svg, fig_4_6_cache_hierarchy.svg)"],
    "MMU dịch địa chỉ ảo sang vật lý; chưa map thì rơi về identity (VA=PA) để chương trình bare-metal vẫn chạy. TLB 8 "
    "entry fully-associative LRU, kiểm quyền R/W/X/cacheable; map/unmap qua syscall. Cache L1 lệnh/dữ liệu 16 set 4 way, "
    "L2 64 set 4 way, block 64B, write-through. MMIO non-cacheable, bypass cache. BẪY: test có tlbWays=2 để kiểm LRU "
    "nhưng cấu hình SoC THẬT hardcode fully-associative — class tổng quát hơn để kiểm chứng được, production là fully-associative.")

add_content(2, "Bus TileLink (UL/UH) & DMA",
    ["TileLink: kênh A (yêu cầu) / D (phản hồi) mức giao dịch; cầu nối UH↔UL hai chiều; burst đa-beat (beat là lũy thừa 2)",
     "DMA: descriptor 3-word (src/dst/config), FIFO đọc và ghi độc lập, copy RAM↔RAM hoặc RAM→ngoại vi",
     "Backpressure (canAcceptTx): khi FIFO TX đầy thì DMA giữ giao dịch, không làm mất byte",
     "(Chèn hình: fig_4_7_tilelink_a_d_frame.svg, fig_4_10_dma_transfer_flow.svg)"],
    "TileLink dùng 2 kênh: A mang yêu cầu, D mang phản hồi — mức giao dịch. Cầu nối UH–UL hai chiều. UH hỗ trợ burst "
    "nhiều beat (beat lũy thừa 2). DMA lập trình qua descriptor 3 từ (nguồn/đích/cấu hình), 2 FIFO đọc-ghi độc lập, copy "
    "RAM–RAM hoặc RAM ra ngoại vi. Điểm xử lý kỹ: backpressure — FIFO phát UART đầy thì DMA giữ giao dịch, không mất byte. "
    "SỐ LIỆU MẠNH: trước backpressure gửi 65 byte chỉ ra ~16; sau khi sửa: đủ 65/65, 0 mất.")

add_content(2, "Ngoại vi & Giao diện web trực quan",
    ["Ngoại vi MMIO: UART (baud theo divisor), LED Matrix 32×32 (VRAM), mouse, keyboard — tất cả đọc bằng polling",
     "Giao diện web: trình soạn thảo + gợi ý lệnh, bảng thanh ghi/bộ nhớ, sơ đồ SoC 'sống', bảng cache/TLB",
     "Bảng log hệ thống lọc theo module: CPU / MMU / cache / TileLink / DMA / I/O",
     "Chạy trực tiếp trên trình duyệt, không cần cài đặt — (chèn ảnh giao diện)"],
    "Các ngoại vi đều ánh xạ bộ nhớ: UART baud qua divisor, LED Matrix 32×32 dạng VRAM, chuột, bàn phím — đọc bằng "
    "polling. Giao diện web: trình soạn thảo có gợi ý lệnh, bảng thanh ghi/bộ nhớ, sơ đồ SoC sống, bảng cache/TLB, và "
    "bảng log lọc theo từng module. Toàn bộ chạy trực tiếp trên trình duyệt. CHÈN ảnh chụp giao diện.")

# ---- PHẦN 03 ----
add_divider("KẾT QUẢ CỦA GIẢI PHÁP ĐỀ XUẤT", "PHẦN 03")

add_image_slide(3, "Phương pháp kiểm chứng — 3 lớp",
    ["Lớp 1 — Kiểm thử đơn vị/khối bằng Node.js cho từng mô-đun (assembler, CPU-F, MMU, cache, TileLink, DMA, UART, log)",
     "Lớp 2 — Đối chiếu MÃ HÓA lệnh với GNU binutils trên bộ chuẩn quốc tế riscv-tests",
     "Lớp 3 — Đối chiếu THỰC THI trên Spike (mô phỏng tham chiếu chính thức của RISC-V)"],
    "Nhóm kiểm chứng theo 3 lớp bổ trợ. Lớp 1: unit test từng mô-đun bằng Node. Lớp 2: đối chiếu MÃ HÓA lệnh với GNU "
    "binutils trên bộ test ISA chính thức riscv-tests. Lớp 3: đối chiếu THỰC THI trên Spike, mô phỏng tham chiếu của "
    "RISC-V. Điểm quan trọng: hai lớp ngoài dùng công cụ và dữ liệu của cộng đồng quốc tế → khách quan, không tự viết "
    "test rồi tự đúng. Hình dưới là dashboard tổng quan kết quả (cả 3 lớp đều xanh).",
    DASH)

add_content(3, "Kết quả mô phỏng (định lượng)",
    ["14/14 script kiểm thử đơn vị PASS — chạy lại trực tiếp 14/06/2026",
     "Đối chiếu GNU binutils: 61 file ELF, 17.978 lệnh khớp, 0 mismatch; 1.456 lệnh CSR/đặc quyền (ngoài phạm vi) skip",
     "Đối chiếu Spike: 61/61 ELF pass (ISA RV32IMF_zicclsm)",
     "Cả 3 lớp tái lập được ngay trên máy demo (script Node tự bắc cầu sang WSL)"],
    "Số liệu: 14/14 script unit đạt, chạy lại 14/06. Đối chiếu GNU binutils trên 61 ELF riscv-tests: 17.978 lệnh trùng "
    "khớp, KHÔNG lệnh nào sai; 1.456 lệnh CSR/đặc quyền ngoài phạm vi được bỏ qua có thống kê. Đối chiếu Spike: 61/61 ELF "
    "đạt. Cả 3 lớp tái lập ngay trên máy demo — script Node tự gọi sang WSL chạy GNU/Spike. CÂU CHỐT MẠNH: em có thể chạy "
    "lại tại chỗ nếu thầy/cô muốn (node test/verify_riscv_tests_spike.mjs).")

add_content(3, "Kết quả hiển thị / Demo thực tế",
    ["LED Matrix: tô gradient 32×32; DMA tô hàng LED trong khi CPU tính toán song song",
     "UART ← DMA: phát đủ 65/65 byte, 0 byte mất nhờ backpressure (fill → stall → drain)",
     "MMU: bảng TLB cập nhật trực tiếp, VA 0x5020 dịch sang PA 0x10010020",
     "Sơ đồ SoC 'sống': nhiều đường bus active đồng thời, badge đếm giao dịch tăng — (chèn ảnh/video demo)"],
    "Kết quả hiển thị thực tế: chương trình tô gradient lên LED 32×32; DMA tự tô một hàng LED trong khi CPU tính song "
    "song. UART: DMA đẩy 65 byte, backpressure giữ đủ 65/65, 0 mất. MMU: bảng TLB cập nhật ngay, VA 0x5020 → PA "
    "0x10010020. Sơ đồ SoC: nhiều đường bus sáng đồng thời, badge đếm giao dịch tăng. MẸO DEMO: chạy qua tools/dev_server.py "
    "(no-cache); demo vòng lặp vô hạn thì kết thúc bằng Stop; có video dự phòng. CHÈN ảnh/video demo.")

# ---- PHẦN 04 ----
add_divider("KẾT LUẬN VÀ HƯỚNG PHÁT TRIỂN", "PHẦN 04")

add_content(4, "Kết luận & Đóng góp",
    ["Hoàn thiện công cụ mô phỏng SoC RISC-V chạy trên web, trực quan hóa toàn bộ datapath",
     "Kiểm chứng 3 lớp với công cụ chuẩn quốc tế (GNU binutils + Spike) — 0 sai",
     "Đóng góp: công cụ dạy/học Kiến trúc máy tính & Hệ điều hành; mã nguồn mở; kỷ luật về tính trung thực phạm vi",
     "Hạn chế trung thực: chưa có IRQ/CSR/đặc quyền; TileLink chưa coherence; FPU chưa đủ rounding-mode"],
    "Kết luận: hoàn thiện công cụ mô phỏng SoC RISC-V chạy web, trực quan toàn datapath, kiểm chứng 3 lớp với công cụ "
    "chuẩn quốc tế — 0 sai về mã hóa lẫn thực thi trong phạm vi. Đóng góp: công cụ dạy/học KTMT & HĐH, mã nguồn mở. "
    "Hạn chế nêu THẲNG: chưa có ngắt/CSR/đặc quyền; bus chưa coherence; FPU chưa đủ rounding-mode. Chủ động nêu hạn chế "
    "> bị động — giảng viên phản biện đánh giá cao sự thành thật.")

add_content(4, "Hướng phát triển",
    ["Bổ sung ngắt (IRQ) và chế độ đặc quyền/CSR (satp, trap, page-walk phần cứng)",
     "Mở rộng TileLink: kênh B/C/E, coherence TL-C, phân xử đa master",
     "Hoàn thiện FPU (FCSR/fflags/frm); mở rộng CAN (CRC, ACK, arbitration bit-level)",
     "Trực quan hóa pipeline; bổ sung benchmark định lượng hiệu năng"],
    "Hướng phát triển: thêm ngắt và chế độ đặc quyền/CSR; mở rộng TileLink (B/C/E, coherence, đa master); hoàn thiện FPU "
    "và mở rộng CAN; trực quan hóa pipeline và bổ sung benchmark định lượng.")

add_content(4, "Phân công công việc & Tài liệu tham khảo",
    ["Phân công Lộc / Khang theo từng mô-đun và giai đoạn (Bảng 1.3) — phục vụ tiêu chí lập kế hoạch/quản lý (LO7)",
     "Tiến độ thực hiện theo các mốc thời gian",
     "Tài liệu tham khảo chính: RISC-V ISA Spec, TileLink Spec, riscv-tests, Spike, M_CAN manual, GNU binutils"],
    "Phân công: [tên] phụ trách [các phần], [tên] phụ trách [các phần], theo tiến độ từng giai đoạn (Bảng 1.3). Tài liệu "
    "tham khảo: đặc tả RISC-V ISA, TileLink, riscv-tests, Spike, M_CAN. SLIDE NÀY phục vụ tiêu chí LO7 — đừng bỏ. CẦN ĐIỀN "
    "bảng phân công thật.")

# ---- PHẦN 05 ----
add_divider("HỎI & ĐÁP", "PHẦN 05",
    "Phần trình bày đến đây là hết. Em cảm ơn thầy/cô đã lắng nghe và rất mong nhận được câu hỏi, góp ý để hoàn thiện "
    "trước buổi bảo vệ.")

# Phụ lục
add_content(5, "Phụ lục — Phòng thủ câu hỏi phản biện",
    ["Vì sao chọn TileLink mà không phải AXI/Wishbone?",
     "Vì sao dùng polling mà không có ngắt? (đúng phạm vi đề tài)",
     "fcvt.w.s(6.5)=7 vì chưa có rounding-mode/FCSR — giải thích trung thực",
     "a0≠0 ở vài demo là DỮ LIỆU (tổng/giá trị), không phải mã lỗi",
     "Class MMU hỗ trợ set/way nhưng cấu hình SoC hardcode fully-associative",
     "Vì sao đối chiếu GNU/Spike là đáng tin (bộ chuẩn cộng đồng)"],
    "TileLink vs AXI/Wishbone: mở, phân tầng (UL đơn giản/UH burst), kênh A/D đủ minh họa giáo dục, không phức tạp như "
    "AXI 5 kênh. Polling vì chưa hiện thực IRQ/CSR — lựa chọn thiết kế đã nêu, nằm trong hướng phát triển. fcvt 6.5→7 vì "
    "chưa có FCSR/frm (mặc định làm tròn) — hạn chế đã biết. a0≠0 ở dma_demo/bus_demo/mmu_syscall_test/test_cache là DỮ "
    "LIỆU, không phải lỗi. MMU class hỗ trợ set/way để kiểm LRU, SoC production fully-associative. GNU/Spike đáng tin vì là "
    "công cụ + dữ liệu chuẩn của cộng đồng quốc tế. (Chi tiết đầy đủ 12 câu trong docs/kltn/ghi_chu_slide_phan_bien.md.)")

prs.save(OUT)
print("Saved:", OUT)
print("Total slides:", len(prs.slides._sldIdLst))
